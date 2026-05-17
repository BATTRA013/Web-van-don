from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

path = Path(r'd:\laragon\www\Web-van-don\docs\DTH225665_LamThanhHuu.pptx')
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

NAVY = RGBColor(11, 20, 38)
BLUE = RGBColor(37, 99, 235)
SKY = RGBColor(14, 165, 233)
GOLD = RGBColor(245, 158, 11)
GREEN = RGBColor(22, 163, 74)
RED = RGBColor(220, 38, 38)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
MID = RGBColor(226, 232, 240)
DARK = RGBColor(30, 41, 59)
MUTED = RGBColor(71, 85, 105)
FONT = 'Arial'


def bg(slide, color=LIGHT):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def header(slide, title, subtitle=''):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.72))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.1), Inches(11.4), Inches(0.52))
    tf = tx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph(); p2.space_before = Pt(2)
        r2 = p2.add_run(); r2.text = subtitle; r2.font.name = FONT; r2.font.size = Pt(10.5); r2.font.color.rgb = MID
    tag = slide.shapes.add_textbox(Inches(10.6), Inches(0.12), Inches(2.1), Inches(0.35))
    ttf = tag.text_frame; ttf.clear()
    p = ttf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = 'Web-van-don'; r.font.name = FONT; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = MID


def footer(slide, text='Báo cáo dự án | Web Quản Lý Vận Đơn'):
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(7.02), Inches(12.2), Inches(0.22))
    tf = tx.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(9); r.font.color.rgb = MUTED


def card(slide, x, y, w, h, color=WHITE, line=MID):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.color.rgb = line
    return shp


def bullets(tf, items, size=13, color=DARK):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.bullet = True; p.space_after = Pt(8)
        r = p.add_run(); r.text = item; r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color


def title_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    for x, y, w, h, c in [(10.7, 0.45, 2.6, 2.6, BLUE), (9.55, 1.25, 1.8, 1.8, SKY), (11.2, 2.05, 1.0, 1.0, GOLD)]:
        shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = c; shp.line.fill.background(); shp.fill.transparency = 0.2
    tx = s.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(9.2), Inches(2.0))
    tf = tx.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = 'Website Quản Lý Vận Đơn Liên Kết Các Hãng Vận Chuyển'; r.font.name = FONT; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    r2 = p2.add_run(); r2.text = 'Báo cáo dự án hệ thống quản lý đơn hàng, gửi chành xe và đồng bộ vận đơn'; r2.font.name = FONT; r2.font.size = Pt(15.5); r2.font.color.rgb = MID
    info = s.shapes.add_textbox(Inches(0.78), Inches(3.05), Inches(6.6), Inches(2.1))
    itf = info.text_frame; itf.clear(); itf.word_wrap = True
    for i, line in enumerate(['Sinh viên: Lâm Thanh Hữu', 'Lớp / Mã SV: DH23TH1, DTH225665', 'Giảng viên hướng dẫn: Huỳnh Lý Thanh Nhàn', 'Ngày báo cáo: 13/05/2026']):
        p = itf.paragraphs[0] if i == 0 else itf.add_paragraph(); p.space_after = Pt(5)
        r = p.add_run(); r.text = line; r.font.name = FONT; r.font.size = Pt(12.5); r.font.color.rgb = WHITE
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(5.55), Inches(2.8), Inches(0.6))
    badge.fill.solid(); badge.fill.fore_color.rgb = GOLD; badge.line.fill.background()
    btf = badge.text_frame; btf.clear(); btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = btf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'PPT báo cáo ngắn gọn'; r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = NAVY
    return s


def two_column_slide(title, subtitle, left_title, left_items, right_title=None, right_items=None, callout=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, title, subtitle)
    card(s, 0.65, 1.1, 6.1, 5.6)
    lt = s.shapes.add_textbox(Inches(0.92), Inches(1.35), Inches(5.5), Inches(0.35))
    ltf = lt.text_frame; ltf.clear()
    p = ltf.paragraphs[0]; r = p.add_run(); r.text = left_title; r.font.name = FONT; r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = BLUE
    lbox = s.shapes.add_textbox(Inches(0.92), Inches(1.82), Inches(5.45), Inches(4.45))
    ltf = lbox.text_frame; ltf.clear(); ltf.word_wrap = True
    bullets(ltf, left_items, size=13)
    if right_title:
        card(s, 6.95, 1.1, 5.75, 5.6)
        rt = s.shapes.add_textbox(Inches(7.18), Inches(1.35), Inches(5.1), Inches(0.35))
        rtf = rt.text_frame; rtf.clear()
        p = rtf.paragraphs[0]; rr = p.add_run(); rr.text = right_title; rr.font.name = FONT; rr.font.size = Pt(17); rr.font.bold = True; rr.font.color.rgb = NAVY
        rbox = s.shapes.add_textbox(Inches(7.18), Inches(1.82), Inches(5.1), Inches(4.45))
        rtf = rbox.text_frame; rtf.clear(); rtf.word_wrap = True
        bullets(rtf, right_items or [], size=13)
    if callout:
        c = card(s, 0.95, 6.0, 11.45, 0.52, color=LIGHT, line=SKY)
        tf = c.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = callout; r.font.name = FONT; r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = NAVY
    footer(s)
    return s


def db_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Xây dựng cơ sở dữ liệu', 'Các bảng nghiệp vụ và thông tin được lưu trong từng bảng')
    card(s, 0.6, 1.1, 12.1, 5.55)
    groups = [
        ('nguoi_dung', BLUE, 0.9, 1.45, 'Lưu hồ sơ tài khoản đăng nhập, vai trò sử dụng hệ thống và trạng thái duyệt tài khoản.'),
        ('don_hang', SKY, 0.9, 2.42, 'Lưu thông tin đơn hàng chính: người nhận, địa chỉ giao, khối lượng, COD, phí ship và trạng thái xử lý.'),
        ('hang_van_chuyen', GOLD, 0.9, 3.39, 'Lưu cấu hình kết nối hãng vận chuyển như token, shop_id, môi trường và owner cấu hình.'),
        ('chi_tiet_don_hang', GREEN, 0.9, 4.36, 'Lưu danh sách sản phẩm theo từng đơn: tên sản phẩm, số lượng, giá và khối lượng.'),
        ('doi_soat_cod', RED, 0.9, 5.33, 'Lưu dữ liệu đối soát COD giữa kỳ vọng và thực nhận để kiểm tra chênh lệch theo đơn.'),
        ('nha_xe', BLUE, 6.7, 1.94, 'Lưu thông tin nhà xe đối tác cho luồng gửi ngoài tuyến: tên nhà xe, liên hệ, tuyến đường.'),
        ('van_don_ngoai_tuyen', SKY, 6.7, 3.55, 'Lưu vận đơn ngoài tuyến gắn với đơn hàng và nhà xe, gồm mã biên lai và ảnh biên lai.'),
    ]
    for title, color, x, y, text in groups:
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.65), Inches(0.52))
        tag.fill.solid(); tag.fill.fore_color.rgb = color; tag.line.fill.background()
        tf = tag.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
        w = 2.85 if x > 5 else 2.95
        desc = s.shapes.add_textbox(Inches(x), Inches(y + 0.56), Inches(w), Inches(1.0))
        dtf = desc.text_frame; dtf.clear(); dtf.word_wrap = True
        p = dtf.paragraphs[0]
        r = p.add_run(); r.text = text; r.font.name = FONT; r.font.size = Pt(10.3); r.font.color.rgb = DARK
    note = s.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11.2), Inches(0.28))
    ntf = note.text_frame; ntf.clear()
    p = ntf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'Trọng tâm trình bày là ý nghĩa nghiệp vụ của bảng thay vì liệt kê toàn bộ cột.'; r.font.name = FONT; r.font.size = Pt(10.2); r.font.italic = True; r.font.color.rgb = MUTED
    footer(s)
    return s


def tools_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Công cụ xây dựng web', 'Các công nghệ chính được dùng trong dự án')
    cards = [
        ('Laravel / PHP', BLUE, ['Xây dựng controller, request, service', 'Mô hình MVC rõ ràng', 'Dễ tổ chức nghiệp vụ vận đơn']),
        ('MySQL', SKY, ['Lưu dữ liệu đơn hàng và carrier', 'Quan hệ bảng chặt chẽ', 'Phù hợp dữ liệu nghiệp vụ']),
        ('Blade / Frontend', GOLD, ['Hiển thị form và bảng dữ liệu', 'Tối ưu thao tác người dùng', 'Kết hợp nhanh với Laravel']),
        ('GHN / Viettel API', GREEN, ['Tạo vận đơn', 'Đồng bộ trạng thái', 'Làm việc đa hãng vận chuyển']),
    ]
    xs = [0.75, 3.85, 6.95, 10.05]
    for (title, color, items), x in zip(cards, xs):
        card(s, x, 1.45, 2.85, 4.8)
        top = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.14), Inches(1.62), Inches(2.57), Inches(0.55))
        top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()
        tf = top.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
        box = s.shapes.add_textbox(Inches(x+0.2), Inches(2.25), Inches(2.45), Inches(3.65))
        tf = box.text_frame; tf.clear(); tf.word_wrap = True
        bullets(tf, items, size=11.5)
    footer(s)
    return s


def features_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Một số chức năng tiêu biểu', 'Chỉ nêu vài chức năng trọng tâm để báo cáo và demo')
    cards = [
        ('Quản lý đơn hàng', BLUE, ['Tạo, xem, sửa, xóa đơn', 'Lọc theo trạng thái và mã đơn', 'Gắn chi tiết sản phẩm']),
        ('Gửi chành xe', SKY, ['Tạo yêu cầu gửi đơn', 'Nhận / từ chối / cập nhật biên lai', 'Theo dõi trạng thái bill']),
        ('Tạo vận đơn & đồng bộ', GOLD, ['Tạo GHN hoặc Viettel', 'Lưu tracking code', 'Đồng bộ trạng thái đơn']),
    ]
    xs = [0.75, 4.3, 7.85]
    for (title, color, items), x in zip(cards, xs):
        card(s, x, 1.45, 3.1, 4.8)
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.14), Inches(1.62), Inches(2.82), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        tf = bar.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
        box = s.shapes.add_textbox(Inches(x+0.18), Inches(2.25), Inches(2.7), Inches(3.5))
        tf = box.text_frame; tf.clear(); tf.word_wrap = True
        bullets(tf, items, size=12)
    footer(s)
    return s


def architecture_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Kiến trúc tổng thể', 'Cách các lớp trong hệ thống phối hợp với nhau')
    card(s, 0.6, 1.25, 12.1, 4.95)
    blocks = [
        ('Người dùng', BLUE, 0.95, ['Chủ shop', 'Quản lý chành xe', 'Admin hệ thống']),
        ('Laravel backend', SKY, 3.9, ['Controller', 'Request validation', 'Service xử lý nghiệp vụ']),
        ('CSDL & API ngoài', GOLD, 6.85, ['MySQL lưu dữ liệu', 'GHN API', 'Viettel Post API']),
        ('Kết quả hiển thị', GREEN, 9.8, ['Dashboard', 'Danh sách đơn', 'Trạng thái vận đơn']),
    ]
    for title, color, x, items in blocks:
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(2.45), Inches(3.1))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
        box = s.shapes.add_textbox(Inches(x + 0.18), Inches(2.42), Inches(2.1), Inches(2.1))
        btf = box.text_frame; btf.clear(); btf.word_wrap = True
        bullets(btf, items, size=10.8)
    footer(s)
    return s


def relations_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Mô hình quan hệ dữ liệu', 'Khóa ngoại và hành vi xóa/cập nhật')
    card(s, 0.7, 1.25, 11.95, 5.0)
    main = [
        ('nguoi_dung', BLUE, 0.9, 1.75),
        ('hang_van_chuyen', SKY, 2.95, 2.05),
        ('don_hang', GOLD, 5.25, 1.55),
        ('chi_tiet_don_hang', GREEN, 7.1, 2.1),
        ('doi_soat_cod', RED, 9.45, 1.6),
        ('nha_xe', BLUE, 11.35, 1.1),
        ('van_don_ngoai_tuyen', SKY, 10.45, 2.5),
    ]
    for name, color, x, w in main:
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(w), Inches(0.8))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = name; r.font.name = FONT; r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = WHITE
    arrows = s.shapes.add_textbox(Inches(0.85), Inches(2.95), Inches(11.7), Inches(0.48))
    atf = arrows.text_frame; atf.clear()
    p = atf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'nguoi_dung 1 - n don_hang | don_hang 1 - n chi_tiet_don_hang | don_hang 1 - n doi_soat_cod'; r.font.name = FONT; r.font.size = Pt(10.6); r.font.bold = True; r.font.color.rgb = DARK
    p2 = atf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = 'don_hang 1 - n van_don_ngoai_tuyen | nha_xe 1 - n van_don_ngoai_tuyen | nguoi_dung 1 - n hang_van_chuyen (nullable)'; r2.font.name = FONT; r2.font.size = Pt(10.3); r2.font.bold = True; r2.font.color.rgb = DARK
    left = s.shapes.add_textbox(Inches(0.95), Inches(3.6), Inches(5.6), Inches(1.9))
    ltf = left.text_frame; ltf.clear(); ltf.word_wrap = True
    bullets(ltf, [
        'Hệ thống giữ toàn vẹn dữ liệu bằng cách ràng buộc chặt bảng don_hang với người dùng và hãng vận chuyển.',
        'Khi xóa đơn hàng, các bản ghi chi tiết đơn, đối soát và vận đơn ngoài tuyến liên quan sẽ tự dọn theo.',
        'Cấu hình hãng vận chuyển vẫn được giữ nếu owner bị xóa (gán nullable để không mất cấu hình API).'
    ], size=11.5)
    right = s.shapes.add_textbox(Inches(6.55), Inches(3.6), Inches(5.5), Inches(1.9))
    rtf = right.text_frame; rtf.clear(); rtf.word_wrap = True
    bullets(rtf, [
        'Luồng ngoài tuyến được tách riêng qua cặp bảng nha_xe và van_don_ngoai_tuyen để dễ kiểm soát vận chuyển thủ công.',
        'Mỗi nghiệp vụ vẫn có mã định danh duy nhất để tránh trùng dữ liệu khi tra cứu hoặc đối soát.',
        'Cách tách bảng này giúp mở rộng thêm báo cáo hoặc tích hợp nhà xe mới mà không ảnh hưởng luồng chính.'
    ], size=11.5)
    footer(s)
    return s


def order_management_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Chức năng quản lý đơn hàng', 'Luồng chính của nghiệp vụ nội bộ')
    card(s, 0.7, 1.3, 12.0, 4.9)
    left = s.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.5), Inches(4.0))
    tf = left.text_frame; tf.clear(); tf.word_wrap = True
    bullets(tf, [
        'Tạo đơn mới với thông tin người nhận, hàng hóa và ghi chú.',
        'Cập nhật đơn theo trạng thái xử lý thực tế.',
        'Tra cứu, lọc và xem chi tiết theo mã đơn hoặc trạng thái.'
    ], size=13)
    right = card(s, 6.95, 1.7, 5.2, 3.7, color=LIGHT, line=BLUE)
    rtf = right.text_frame; rtf.clear(); rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = rtf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'Tập trung vào CRUD + trạng thái'; r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BLUE
    p2 = rtf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = 'Mục tiêu là giảm thao tác thủ công và giúp quản lý đơn rõ ràng hơn.'; r2.font.name = FONT; r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED
    footer(s)
    return s


def route_bill_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Chức năng gửi chành xe', 'Xử lý đơn chuyển qua nhà xe / chành xe')
    card(s, 0.7, 1.3, 12.0, 4.9)
    cols = [
        ('Tạo yêu cầu', SKY, ['Chọn đơn cần gửi', 'Sinh mã theo dõi', 'Lưu thông tin chuyến gửi']),
        ('Xử lý biên lai', GOLD, ['Nhận biên lai hoặc ảnh minh chứng', 'Cập nhật receipt', 'Đối soát kết quả']),
        ('Trạng thái', GREEN, ['Chờ duyệt', 'Đã nhận', 'Từ chối / hoàn tất']),
    ]
    xs = [1.0, 4.55, 8.1]
    for (title, color, items), x in zip(cols, xs):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.75), Inches(2.9))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
        box = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.55), Inches(2.4), Inches(2.0))
        btf = box.text_frame; btf.clear(); btf.word_wrap = True
        bullets(btf, items, size=11)
    footer(s)
    return s


def shipment_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Chức năng tạo vận đơn', 'Làm việc với GHN và Viettel Post')
    card(s, 0.7, 1.3, 12.0, 4.9)
    left = s.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.5), Inches(4.0))
    tf = left.text_frame; tf.clear(); tf.word_wrap = True
    bullets(tf, [
        'Nhập dữ liệu gửi hàng và kiểm tra thông tin người nhận.',
        'Chuẩn bị payload theo từng hãng vận chuyển.',
        'Nhận tracking code và lưu vào hệ thống để tra cứu về sau.'
    ], size=13)
    right = card(s, 6.95, 1.7, 5.2, 3.7, color=LIGHT, line=GOLD)
    rtf = right.text_frame; rtf.clear(); rtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = rtf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'GHN / Viettel Post'; r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = GOLD
    p2 = rtf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)
    r2 = p2.add_run(); r2.text = 'Hệ thống dùng chung một luồng, chỉ khác phần cấu hình và mapping dữ liệu.'; r2.font.name = FONT; r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED
    footer(s)
    return s


def sync_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Chức năng đồng bộ trạng thái', 'Cập nhật trạng thái từ hãng vận chuyển về hệ thống')
    card(s, 0.7, 1.3, 12.0, 4.9)
    left = s.shapes.add_textbox(Inches(1.0), Inches(1.65), Inches(5.6), Inches(4.0))
    tf = left.text_frame; tf.clear(); tf.word_wrap = True
    bullets(tf, [
        'Đồng bộ theo từng đơn khi cần xem lại trạng thái cụ thể.',
        'Có thể chạy hàng loạt cho nhiều đơn cùng lúc.',
        'Sau khi nhận kết quả, hệ thống map về trạng thái nội bộ thống nhất.'
    ], size=13)
    right = s.shapes.add_textbox(Inches(6.95), Inches(1.9), Inches(4.9), Inches(2.9))
    rtf = right.text_frame; rtf.clear(); rtf.word_wrap = True
    bullets(rtf, [
        'Tránh lệch trạng thái giữa hệ thống và hãng vận chuyển.',
        'Giảm thao tác thủ công cho người quản lý.'
    ], size=12.5)
    footer(s)
    return s


def permission_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Phân quyền và an toàn dữ liệu', 'Mỗi vai trò chỉ thao tác trong phạm vi được phép')
    card(s, 0.7, 1.3, 12.0, 4.9)
    roles = [
        ('Chủ shop', BLUE, ['Tạo và quản lý đơn của mình', 'Theo dõi vận đơn', 'Xem kết quả xử lý']),
        ('Quản lý chành xe', SKY, ['Nhận / từ chối đơn gửi', 'Cập nhật biên lai', 'Xem trạng thái liên quan']),
        ('Admin', GOLD, ['Quản trị dữ liệu toàn hệ thống', 'Xem tổng quan hoạt động', 'Kiểm soát quyền truy cập']),
    ]
    xs = [1.0, 4.55, 8.1]
    for (title, color, items), x in zip(roles, xs):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.75), Inches(2.9))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
        box = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.55), Inches(2.4), Inches(2.0))
        btf = box.text_frame; btf.clear(); btf.word_wrap = True
        bullets(btf, items, size=11)
    footer(s)
    return s


def status_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Trạng thái và mapping nghiệp vụ', 'Chuẩn hóa trạng thái từ nhiều nguồn về một hệ thống')
    card(s, 0.75, 1.35, 12.0, 4.9)
    left = s.shapes.add_textbox(Inches(1.05), Inches(1.7), Inches(5.4), Inches(4.0))
    tf = left.text_frame; tf.clear(); tf.word_wrap = True
    bullets(tf, [
        'Mỗi hãng có tên trạng thái riêng nên cần chuẩn hóa trước khi hiển thị.',
        'Hệ thống đổi về bộ trạng thái nội bộ để báo cáo thống nhất.',
        'Nhờ đó, người dùng chỉ cần nhìn một cách hiểu duy nhất.'
    ], size=13)
    right = s.shapes.add_textbox(Inches(6.95), Inches(1.7), Inches(4.9), Inches(3.8))
    rtf = right.text_frame; rtf.clear(); rtf.word_wrap = True
    bullets(rtf, [
        'Trạng thái nhận được từ GHN / Viettel.',
        'Bộ map nội bộ của hệ thống.',
        'Trạng thái hiển thị trên dashboard.'
    ], size=12.5)
    footer(s)
    return s


def summary_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Một số chức năng tiêu biểu', 'Tóm tắt nhanh các nhóm chức năng chính')
    card(s, 0.75, 1.45, 12.0, 4.7)
    cols = [
        ('Đơn hàng', BLUE, ['CRUD đơn', 'Chi tiết đơn', 'Lọc theo trạng thái']),
        ('Vận chuyển', SKY, ['Gửi chành xe', 'Tạo vận đơn', 'Đồng bộ trạng thái']),
        ('Quản trị', GOLD, ['Phân quyền', 'Theo dõi dữ liệu', 'Kiểm soát truy cập']),
    ]
    xs = [1.05, 4.55, 8.05]
    for (title, color, items), x in zip(cols, xs):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.95), Inches(2.6), Inches(3.0))
        shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
        tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.name = FONT; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
        box = s.shapes.add_textbox(Inches(x + 0.15), Inches(2.5), Inches(2.3), Inches(2.0))
        btf = box.text_frame; btf.clear(); btf.word_wrap = True
        bullets(btf, items, size=11)
    footer(s)
    return s


def demo_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Kịch bản demo trực tiếp', 'Phần này trình bày bằng thao tác trên web cho giáo viên xem')
    card(s, 0.8, 1.35, 11.75, 5.0)
    tx = s.shapes.add_textbox(Inches(1.15), Inches(1.65), Inches(5.6), Inches(4.1))
    tf = tx.text_frame; tf.clear(); tf.word_wrap = True
    bullets(tf, [
        'Đăng nhập vào hệ thống với tài khoản đã chuẩn bị sẵn.',
        'Mở một đơn hàng mẫu để giáo viên xem dữ liệu thực tế.',
        'Thử tạo vận đơn GHN hoặc Viettel trên form.',
        'Bấm đồng bộ trạng thái để thấy dữ liệu cập nhật ngay.',
        'Giải thích ngắn gọn vai trò từng màn hình, không đi sâu code.'
    ], size=14)
    box = card(s, 7.1, 1.7, 4.9, 3.9, color=LIGHT, line=BLUE)
    tf2 = box.text_frame; tf2.clear(); tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'Chỗ chèn ảnh / màn hình demo'; r.font.name = FONT; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BLUE
    p2 = tf2.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = 'Tài khoản test, danh sách đơn, form tạo vận đơn...'; r2.font.name = FONT; r2.font.size = Pt(11.5); r2.font.color.rgb = MUTED
    footer(s)
    return s


def result_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, 'Kết quả và hướng phát triển', 'Kết lại phần trình bày trước khi demo')
    card(s, 0.75, 1.35, 12.0, 4.9)
    left = s.shapes.add_textbox(Inches(1.05), Inches(1.7), Inches(5.5), Inches(4.0))
    tf = left.text_frame; tf.clear(); tf.word_wrap = True
    bullets(tf, [
        'Hệ thống đã gom được các luồng đơn hàng, chành xe và vận đơn về một nơi.',
        'Quy trình xử lý rõ ràng hơn, giảm thao tác thủ công và hạn chế sai lệch dữ liệu.',
        'Phân quyền theo vai trò giúp đảm bảo mỗi người chỉ thao tác trong phạm vi của mình.'
    ], size=13.2)
    right = s.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.2), Inches(3.9))
    tf2 = right.text_frame; tf2.clear(); tf2.word_wrap = True
    bullets(tf2, [
        'Bổ sung dashboard thống kê và biểu đồ.',
        'Mở rộng thêm carrier nếu phát sinh nhu cầu.',
        'Tối ưu giao diện và log lịch sử thao tác.',
        'Hoàn thiện thêm các luồng kiểm tra / đối soát.'
    ], size=13.2)
    footer(s)
    return s


def closing_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, NAVY)
    panel = card(s, 1.0, 1.3, 11.35, 4.8, color=WHITE, line=NAVY)
    tf = panel.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'Kết luận'; r.font.name = FONT; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(14)
    r2 = p2.add_run(); r2.text = 'Đề tài giúp quản lý vận đơn tập trung, có thể demo trực tiếp trên web để làm rõ chức năng.'; r2.font.name = FONT; r2.font.size = Pt(17); r2.font.color.rgb = DARK
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(12)
    r3 = p3.add_run(); r3.text = 'Em xin cảm ơn thầy/cô đã lắng nghe.'; r3.font.name = FONT; r3.font.size = Pt(15); r3.font.bold = True; r3.font.color.rgb = BLUE
    footer(s, '')
    return s

# Build deck

title_slide()
two_column_slide('Các vấn đề cấp thiết', 'Vì sao em chọn đề tài này', 'Thực trạng cần giải quyết', [
    'Đơn hàng bị phân tán ở nhiều nơi, khó theo dõi tập trung.',
    'Mỗi hãng vận chuyển có quy trình và trạng thái khác nhau.',
    'Nếu xử lý thủ công dễ sai lệch tracking và chậm cập nhật trạng thái.',
    'Cần hệ thống quản lý thống nhất, dễ mở rộng và dễ kiểm soát quyền.'
], 'Ý nghĩa của đề tài', [
    'Tập trung dữ liệu vào một hệ thống.',
    'Chuẩn hóa luồng xử lý đơn hàng.',
    'Giảm sai sót khi làm việc với nhiều carrier.'
], 'Bài toán chính: tập trung dữ liệu và chuẩn hóa luồng xử lý.')
two_column_slide('Mục tiêu và phạm vi', 'Hệ thống tập trung vào những gì', 'Mục tiêu chính', [
    'Quản lý đơn hàng nội bộ.',
    'Gửi chành xe và cập nhật biên lai.',
    'Tạo vận đơn GHN / Viettel Post.',
    'Đồng bộ trạng thái đơn.',
    'Phân quyền theo vai trò người dùng.'
], 'Phạm vi sử dụng', [
    'Chủ shop: tạo và theo dõi đơn của mình.',
    'Quản lý chành xe: nhận / từ chối / cập nhật bill.',
    'Admin: quản trị dữ liệu và theo dõi toàn hệ thống.'
], 'Chỉ chọn các chức năng cốt lõi để báo cáo ngắn gọn.')
tools_slide()
architecture_slide()
db_slide()
relations_slide()
order_management_slide()
route_bill_slide()
shipment_slide()
sync_slide()
permission_slide()
status_slide()
result_slide()
closing_slide()

prs.save(str(path))
print('saved', path)
print('slides', len(prs.slides))
